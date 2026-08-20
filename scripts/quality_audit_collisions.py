import sys, io, os, re
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import extract_financial_ppt as ex
from pptx import Presentation
from collections import defaultdict

BASE_DIR = ex.BASE_DIR
TITLE_KEYWORD = ex.TITLE_KEYWORD

_STAGE_SUFFIXES = ["사전검토", "착수", "중간", "완료", "제안"]


def strip_stage_suffix(filename):
    name = os.path.splitext(filename)[0]
    for suf in _STAGE_SUFFIXES:
        name = re.sub(rf"[_\[\(]?{re.escape(suf)}[\]\)]?$", "", name).strip()
    return name


def read_header(tbl):
    headers = {}
    for c in range(len(tbl.columns)):
        try:
            headers[c] = tbl.cell(0, c).text.strip()
        except Exception:
            headers[c] = ""
    return headers


def find_col(headers, keywords):
    for idx, txt in headers.items():
        if any(kw in txt for kw in keywords):
            return idx
    return -1


groups = defaultdict(list)
aip_files = []
open_errors = []
short_col_files = []

for root, _, files in os.walk(BASE_DIR):
    for name in files:
        if name.startswith("~$"):
            continue
        ext = os.path.splitext(name)[1].lower()
        if ext != ".pptx":
            continue
        full_path = os.path.join(root, name)

        if ex.is_aip_encrypted(full_path):
            aip_files.append(full_path)
            continue

        try:
            prs = Presentation(full_path)
        except Exception as e:
            open_errors.append((full_path, str(e)))
            continue

        year = ex.extract_year_from_filename(name)
        part = ex.extract_part_from_path(full_path)

        for slide in prs.slides:
            has_kw = False
            for shape in slide.shapes:
                if shape.has_text_frame and TITLE_KEYWORD in shape.text_frame.text:
                    has_kw = True
                    break
            if not has_kw:
                continue

            best = None
            best_cols = 0
            for shape in slide.shapes:
                if shape.has_table:
                    t = shape.table
                    if len(t.columns) > best_cols:
                        best_cols = len(t.columns)
                        best = t
            if best is None:
                continue

            tbl = best
            ncols = len(tbl.columns)
            if ncols < 10:
                short_col_files.append((full_path, ncols))
                continue

            headers = read_header(tbl)
            has_jaegyungbi = ncols >= 11
            offset = 1 if has_jaegyungbi else 0
            idx_note = find_col(headers, ["비고"])
            if idx_note < 0:
                idx_note = 9 + offset

            for r in range(1, len(tbl.rows)):
                try:
                    code = ex.safe_str(tbl.cell(r, 0).text)
                    gubun = ex.safe_str(tbl.cell(r, 1).text)
                except Exception:
                    continue
                if code in ("", "0") and gubun in ("", "0"):
                    continue
                try:
                    note_raw = tbl.cell(r, min(idx_note, ncols - 1)).text
                except Exception:
                    note_raw = ""
                note = ex.clean_note_value(note_raw)
                key = (code, str(year), part, gubun)
                groups[key].append((full_path, note))

dupes = {k: v for k, v in groups.items() if len(v) > 1}

same_project_dupes = []
cross_project_dupes = []

for key, entries in dupes.items():
    bases = {strip_stage_suffix(os.path.basename(p)) for p, _ in entries}
    if len(bases) == 1:
        same_project_dupes.append((key, entries))
    else:
        cross_project_dupes.append((key, entries))

print(f"[AIP 암호화로 읽기 불가] {len(aip_files)}개")
for p in aip_files:
    print(f"  - {p}")

print(f"\n[열기 실패] {len(open_errors)}개")
for p, e in open_errors:
    print(f"  - {p} : {e}")

print(f"\n[컬럼 10개 미만(구 템플릿)] {len(short_col_files)}개")
for p, n in short_col_files:
    print(f"  - {os.path.basename(p)} (cols={n})")

print(f"\n총 (코드,연도,파트,구분) 키 그룹: {len(groups)}")
print(f"중복 키 그룹(2개 이상 파일에서 생성): {len(dupes)}")
print(f"  - 같은 프로젝트 단계파일 간 중복(안전, 파일명 베이스 동일): {len(same_project_dupes)}")
print(f"  - 서로 다른 프로젝트 간 코드 충돌(위험!): {len(cross_project_dupes)}")

print("\n" + "=" * 100)
print("[위험] 서로 다른 프로젝트 코드 충돌 목록")
print("=" * 100)
for key, entries in cross_project_dupes:
    print(f"\nKEY = {key}")
    for p, note in entries:
        print(f"  - {os.path.basename(p)}")
        print(f"      비고: {note!r}")

print("\n" + "=" * 100)
print("[참고] 같은 프로젝트 단계파일 간 중복 (착수/완료/제안 파일이 서로 겹치는 stage row를 포함) — 비고 내용 다른 것만 표시")
print("=" * 100)
diff_note_count = 0
for key, entries in same_project_dupes:
    notes = {n for _, n in entries}
    if len(notes) > 1:
        diff_note_count += 1
        print(f"\nKEY = {key}")
        for p, note in entries:
            print(f"  - {os.path.basename(p)}")
            print(f"      비고: {note!r}")
print(f"\n비고 내용이 서로 다른 same-project 중복 그룹 수: {diff_note_count} / 전체 same-project 중복 {len(same_project_dupes)}개 중")
