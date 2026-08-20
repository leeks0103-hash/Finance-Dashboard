# -*- coding: utf-8 -*-
# PPT 파일 전체 품질 감사 분석 스크립트
# 출력: data/품질감사_보고서.xlsx

import os
import sys
import re
import struct
import traceback
import pandas as pd
from openpyxl import Workbook
from openpyxl.styles import PatternFill, Font, Alignment
from openpyxl.utils import get_column_letter

sys.stdout.reconfigure(encoding='utf-8')

# ─── 경로 설정 ────────────────────────────────────────────────────────────────
PPT_ROOT = r'C:\Users\aaa\Desktop\기술교육실_프로젝트 보고서 수집'
EXCEL_PATH = r'C:\Users\aaa\coding\dashboard\data\재무관점 필수 데이터 추출.xlsx'
OUTPUT_PATH = r'C:\Users\aaa\coding\dashboard\data\품질감사_보고서.xlsx'

# ─── 판별 기준 ────────────────────────────────────────────────────────────────
PART_KEYWORDS = ['AI', 'SW', 'PM', '미모', '신사업', '전차', 'K뉴딜TF']
STAGE_KEYWORDS = ['완료', '중간', '착수', '제안', '사전검토', '사업계획', '검토']
# 데이터 누락 판단 대상 단계 (이 단계는 추출 데이터가 있어야 함)
DATA_EXPECTED_STAGES = ['완료', '중간', '진행']
# 데이터가 없어도 정상인 단계
DATA_EXEMPT_STAGES = ['착수', '제안', '사전검토', '사업계획', '검토']

# ─── 유틸 함수 ────────────────────────────────────────────────────────────────
def is_aip(filepath):
    """파일 첫 8바이트로 AIP(CFB) 여부 감지"""
    try:
        with open(filepath, 'rb') as f:
            magic = f.read(8)
        return magic == bytes.fromhex('d0cf11e0a1b11ae1')
    except Exception:
        return False

def contains_any(text, keywords):
    """텍스트에 키워드 중 하나라도 포함되면 True"""
    for kw in keywords:
        if kw in text:
            return True
    return False

def extract_year_from_filename(fname):
    """파일명에서 연도 추출 (25년/26년/25/26 등)"""
    m = re.search(r'(\d{2,4})년', fname)
    if m:
        y = m.group(1)
        return y[-2:]  # 마지막 2자리
    m = re.search(r'_(\d{2})_', fname)
    if m:
        return m.group(1)
    return None

def extract_stage_from_filename(fname):
    """파일명에서 보고단계 추출"""
    for kw in ['완료', '중간', '착수', '제안', '사전검토', '사업계획', '검토']:
        if kw in fname:
            return kw
    return None

def get_slide1_text_pptx(filepath):
    """python-pptx로 슬라이드1 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(filepath)
    if not prs.slides:
        return ''
    slide = prs.slides[0]
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    texts.append(t)
    return ' '.join(texts)

def get_slide1_text_win32(filepath):
    """win32com으로 슬라이드1 텍스트 추출 (AIP 파일용)"""
    import pythoncom
    import win32com.client
    pythoncom.CoInitialize()
    ppt_app = None
    prs = None
    text = ''
    try:
        ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
        ppt_app.Visible = 1
        prs = ppt_app.Presentations.Open(
            os.path.abspath(filepath),
            ReadOnly=True,
            Untitled=False,
            WithWindow=False
        )
        if prs.Slides.Count > 0:
            slide = prs.Slides(1)
            parts = []
            for shape in slide.Shapes:
                try:
                    if shape.HasTextFrame:
                        t = shape.TextFrame.TextRange.Text.strip()
                        if t:
                            parts.append(t)
                except Exception:
                    pass
            text = ' '.join(parts)
    except Exception as e:
        text = f'[오류: {e}]'
    finally:
        try:
            if prs:
                prs.Close()
        except Exception:
            pass
        try:
            if ppt_app:
                ppt_app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()
    return text

# ─── PPT 파일 수집 ─────────────────────────────────────────────────────────────
print("PPT 파일 수집 중...")
ppt_files = []
for root, dirs, files in os.walk(PPT_ROOT):
    # Thumbs.db, 로그, 엑셀 등 제외
    for f in files:
        if f.lower().endswith(('.pptx', '.ppt')) and not f.startswith('~$'):
            folder = os.path.relpath(root, PPT_ROOT)
            fpath = os.path.join(root, f)
            aip = is_aip(fpath)
            ppt_files.append({
                'filepath': fpath,
                'filename': f,
                'folder': folder,
                'aip': '예' if aip else '아니오',
                'aip_bool': aip,
            })
print(f"  총 {len(ppt_files)}개 PPT 파일 발견")

# ─── 엑셀 파일명 목록 로드 ──────────────────────────────────────────────────────
print("엑셀 파일 로드 중...")
df_excel = pd.read_excel(EXCEL_PATH, sheet_name='취합', header=0)
col_fname = df_excel.columns[12]  # 원본파일명
excel_filenames = set(df_excel[col_fname].dropna().astype(str).tolist())
print(f"  엑셀 내 파일명 {len(excel_filenames)}개")

# ═══════════════════════════════════════════════════════════════════════════════
# 분석 1: 파트명 키워드 누락
# ═══════════════════════════════════════════════════════════════════════════════
print("\n[1/4] 파트명 키워드 누락 분석...")
sheet1_rows = []
for p in ppt_files:
    fname = p['filename']
    if not contains_any(fname, PART_KEYWORDS):
        sheet1_rows.append({
            '파일명': fname,
            '폴더': p['folder'],
            'AIP여부': p['aip'],
            '비고': '',
        })
print(f"  누락 파일: {len(sheet1_rows)}개")

# ═══════════════════════════════════════════════════════════════════════════════
# 분석 2: 보고단계 키워드 누락
# ═══════════════════════════════════════════════════════════════════════════════
print("\n[2/4] 보고단계 키워드 누락 분석...")
sheet2_rows = []
for p in ppt_files:
    fname = p['filename']
    if not contains_any(fname, STAGE_KEYWORDS):
        sheet2_rows.append({
            '파일명': fname,
            '폴더': p['folder'],
            'AIP여부': p['aip'],
            '비고': '',
        })
print(f"  누락 파일: {len(sheet2_rows)}개")

# ═══════════════════════════════════════════════════════════════════════════════
# 분석 3: 데이터 누락 파일 (완료/중간/진행인데 엑셀에 없음)
# ═══════════════════════════════════════════════════════════════════════════════
print("\n[3/4] 데이터 누락 파일 분석...")
sheet3_rows = []
for p in ppt_files:
    fname = p['filename']
    stage = extract_stage_from_filename(fname)

    # 제외 단계면 스킵
    if stage in DATA_EXEMPT_STAGES:
        continue
    # 데이터 기대 단계거나 단계 불명인 경우
    if stage in DATA_EXPECTED_STAGES or stage is None:
        # 엑셀에 파일명이 있는지 확인
        if fname not in excel_filenames:
            sheet3_rows.append({
                '파일명': fname,
                '폴더': p['folder'],
                '보고단계': stage if stage else '(미확인)',
                'AIP여부': p['aip'],
                '비고': '',
            })
print(f"  누락 파일: {len(sheet3_rows)}개")

# ═══════════════════════════════════════════════════════════════════════════════
# 분석 4: 파일명 vs 슬라이드1 제목 불일치 (완료/중간 단계만)
# ═══════════════════════════════════════════════════════════════════════════════
print("\n[4/4] 파일명-슬라이드 내용 불일치 분석...")
print("  (완료/중간 단계 파일만 분석)")

target_files = [p for p in ppt_files
                if extract_stage_from_filename(p['filename']) in ['완료', '중간']]
print(f"  대상 파일: {len(target_files)}개")

sheet4_rows = []
for i, p in enumerate(target_files):
    fname = p['filename']
    fpath = p['filepath']
    aip = p['aip_bool']

    print(f"  [{i+1}/{len(target_files)}] {fname[:60]}...", end='', flush=True)

    slide_text = ''
    try:
        if aip:
            slide_text = get_slide1_text_win32(fpath)
        else:
            slide_text = get_slide1_text_pptx(fpath)
        print(' OK')
    except Exception as e:
        slide_text = f'[오류: {e}]'
        print(f' 오류: {e}')

    # 파일명 연도 vs 슬라이드 연도
    fname_year = extract_year_from_filename(fname)
    slide_year = None
    if fname_year:
        # 슬라이드에서 연도 찾기
        m_list = re.findall(r'(\d{2,4})년', slide_text)
        years_in_slide = list(set([y[-2:] for y in m_list]))
        if years_in_slide and fname_year not in years_in_slide:
            slide_year = '/'.join(years_in_slide)
        elif years_in_slide:
            slide_year = '/'.join(years_in_slide)

    # 파일명 단계 vs 슬라이드 단계
    fname_stage = extract_stage_from_filename(fname)
    slide_stages = [kw for kw in STAGE_KEYWORDS if kw in slide_text]
    slide_stage_str = '/'.join(slide_stages) if slide_stages else '(없음)'

    # 불일치 항목 파악
    mismatch_items = []
    year_mismatch = False
    stage_mismatch = False

    if fname_year and slide_year and fname_year not in (slide_year or ''):
        year_mismatch = True
        mismatch_items.append('연도')

    if fname_stage and slide_stages and fname_stage not in slide_stages:
        stage_mismatch = True
        mismatch_items.append('보고단계')

    if mismatch_items or '[오류' in slide_text:
        note = '읽기오류' if '[오류' in slide_text else ''
        sheet4_rows.append({
            '파일명': fname,
            '폴더': p['folder'],
            '파일명연도': fname_year or '',
            '슬라이드연도': slide_year or '',
            '파일명단계': fname_stage or '',
            '슬라이드단계': slide_stage_str,
            '불일치항목': ', '.join(mismatch_items) if mismatch_items else '읽기오류',
            'AIP여부': p['aip'],
        })

print(f"  불일치 파일: {len(sheet4_rows)}개")

# ═══════════════════════════════════════════════════════════════════════════════
# 엑셀 출력
# ═══════════════════════════════════════════════════════════════════════════════
print("\n엑셀 보고서 생성 중...")

HEADER_FILL = PatternFill("solid", fgColor="4472C4")
HEADER_FONT = Font(bold=True, color="FFFFFF")
ALT_FILL = PatternFill("solid", fgColor="F2F2F2")
WHITE_FILL = PatternFill("solid", fgColor="FFFFFF")

def write_sheet(ws, headers, rows, sheet_title=None):
    """시트에 데이터 작성"""
    # 헤더
    ws.append(headers)
    for col_idx, _ in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col_idx)
        cell.fill = HEADER_FILL
        cell.font = HEADER_FONT
        cell.alignment = Alignment(horizontal='center', vertical='center')

    # 데이터 행
    for row_idx, row_data in enumerate(rows, start=2):
        ws.append([row_data.get(h, '') for h in headers])
        fill = ALT_FILL if row_idx % 2 == 0 else WHITE_FILL
        for col_idx in range(1, len(headers) + 1):
            ws.cell(row=row_idx, column=col_idx).fill = fill

    # 열 너비 자동조정
    for col_idx, header in enumerate(headers, 1):
        col_letter = get_column_letter(col_idx)
        max_len = len(header)
        for row_data in rows:
            val = str(row_data.get(header, ''))
            max_len = max(max_len, len(val))
        ws.column_dimensions[col_letter].width = min(max_len + 4, 60)

    # 행 높이
    ws.row_dimensions[1].height = 22

wb = Workbook()
# 기본 시트 제거
wb.remove(wb.active)

# 시트 1: 파트명 누락
ws1 = wb.create_sheet("파트명 누락")
write_sheet(ws1, ['파일명', '폴더', 'AIP여부', '비고'], sheet1_rows)

# 시트 2: 보고단계 누락
ws2 = wb.create_sheet("보고단계 누락")
write_sheet(ws2, ['파일명', '폴더', 'AIP여부', '비고'], sheet2_rows)

# 시트 3: 데이터 누락
ws3 = wb.create_sheet("데이터 누락")
write_sheet(ws3, ['파일명', '폴더', '보고단계', 'AIP여부', '비고'], sheet3_rows)

# 시트 4: 파일명-내용 불일치
ws4 = wb.create_sheet("파일명-내용 불일치")
write_sheet(ws4,
    ['파일명', '폴더', '파일명연도', '슬라이드연도', '파일명단계', '슬라이드단계', '불일치항목', 'AIP여부'],
    sheet4_rows)

# 요약 시트 추가
ws0 = wb.create_sheet("요약", 0)
ws0.column_dimensions['A'].width = 30
ws0.column_dimensions['B'].width = 15
summary_data = [
    ['분석 항목', '건수'],
    ['전체 PPT 파일', len(ppt_files)],
    ['1. 파트명 키워드 누락', len(sheet1_rows)],
    ['2. 보고단계 키워드 누락', len(sheet2_rows)],
    ['3. 데이터 누락 (완료/중간인데 추출 없음)', len(sheet3_rows)],
    ['4. 파일명-슬라이드 내용 불일치', len(sheet4_rows)],
]
for row in summary_data:
    ws0.append(row)
# 헤더 스타일
for cell in ws0[1]:
    cell.fill = HEADER_FILL
    cell.font = HEADER_FONT
    cell.alignment = Alignment(horizontal='center')

wb.save(OUTPUT_PATH)
print(f"\n✓ 저장 완료: {OUTPUT_PATH}")
print("\n=== 결과 요약 ===")
print(f"전체 PPT 파일        : {len(ppt_files)}개")
print(f"1. 파트명 누락       : {len(sheet1_rows)}개")
print(f"2. 보고단계 누락     : {len(sheet2_rows)}개")
print(f"3. 데이터 누락       : {len(sheet3_rows)}개")
print(f"4. 파일명-내용 불일치: {len(sheet4_rows)}개")
