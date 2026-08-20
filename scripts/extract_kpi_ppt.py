import os
import re
import hashlib
import logging
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Optional, List, Dict, Tuple

from dotenv import load_dotenv
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
from openpyxl import load_workbook, Workbook

try:
    import pythoncom
    import win32com.client
    WIN32_AVAILABLE = True
except ImportError:
    WIN32_AVAILABLE = False

load_dotenv()  # .env 파일이 있으면 환경변수로 로드 (없으면 무시)

# =========================================
# 사용자 설정
# =========================================
# 환경변수 EXTRACT_KPI_ROOT_DIR 우선 사용 (.env 또는 compare_and_update.py가 자동 주입)
# CLI 인수로도 덮어쓰기 가능: python extract_kpi_ppt.py "C:\새폴더경로"
import sys as _sys
RETRY_MODE = "--retry" in _sys.argv
ROOT_DIR = Path(os.environ.get(
    "EXTRACT_KPI_ROOT_DIR",
    r"C:\Users\aaa\Desktop\기술교육실_프로젝트 보고서 수집",
))
if not RETRY_MODE and len(_sys.argv) > 1 and Path(_sys.argv[1]).is_dir():
    ROOT_DIR = Path(_sys.argv[1])

# 출력 엑셀: 프로젝트 data/ 폴더로 저장
_DATA_DIR = Path(__file__).resolve().parent.parent / "data"
_DATA_DIR.mkdir(exist_ok=True)
AIP_FAILED_FILE = _DATA_DIR / "kpi_aip_failed.txt"  # AIP 실패 기록

TARGET_EXCEL_NAME = "KPI 지표 데이터 추출.xlsx"
DATA_SHEET_NAME = "취합"
HISTORY_SHEET_NAME = "처리 이력"
SUMMARY_SHEET_NAME = "kpi 집계"

TITLE_KEYWORD = "KPI/경영현황"

PPT_EXTENSIONS = {".ppt", ".pptx"}
START_ROW = 2
END_ROW = 9
KEY_COL_1 = 1
KEY_COL_2 = 2
DATA_COLS = [5, 6, 7, 8]
FLOAT_ROWS = [4, 7, 8]

# 취합 시트 헤더 — KPI 항목 순서 (PPT 표 row 2~9 순서와 동일)
KPI_METRIC_NAMES = [
    "NPS",
    "전략기술과정_건수",
    "전략기술과정_적절성",
    "특화교육체계_건수",
    "AI교육_고객사건수",
    "AI교육_적절성",
    "신사업_매출억",
    "신사업_신규기존건수",
]

PART_KEYWORDS = ["신사업", "PM", "전차", "미모", "AI", "SW", "K뉴딜TF"]
REPORT_STAGE_KEYWORDS = ["사전검토", "사업계획", "제안", "착수", "중간", "완료", "검토"]

# 폴더명 → 파트명 매핑 (파일명에 키워드 없는 경우 상위 폴더명으로 판단)
FOLDER_PART_MAP = {
    "AI교육파트": "AI",
    "SW교육파트": "SW",
    "교육사업PM파트": "PM",
    "신사업기획파트": "신사업",
    "미래모빌리티교육파트": "미모",
    "전동화&차량개발교육파트": "전차",
    "K뉴딜 아카데미 TF": "K뉴딜TF",
}

LOG_FILE_NAME = "extract_kpi_ppt.log"
HASH_CHUNK_SIZE = 1024 * 1024
MAX_SEARCH_DEPTH = 2

ILLEGAL_CHARACTERS_RE = re.compile(r"[\x00-\x08\x0B-\x0C\x0E-\x1F]")


# =========================================
# 로그 설정
# =========================================
def setup_logger(log_dir: Path) -> logging.Logger:
    log_dir.mkdir(parents=True, exist_ok=True)
    log_file_path = log_dir / LOG_FILE_NAME

    logger = logging.getLogger("extract_kpi_ppt")
    logger.setLevel(logging.INFO)

    if logger.handlers:
        return logger

    formatter = logging.Formatter(
        fmt="%(asctime)s | %(levelname)s | %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S"
    )

    file_handler = logging.FileHandler(log_file_path, encoding="utf-8")
    file_handler.setLevel(logging.INFO)
    file_handler.setFormatter(formatter)

    stream_handler = logging.StreamHandler()
    stream_handler.setLevel(logging.INFO)
    stream_handler.setFormatter(formatter)

    logger.addHandler(file_handler)
    logger.addHandler(stream_handler)
    return logger


logger = setup_logger(_DATA_DIR)


# =========================================
# 공통 유틸
# =========================================
def now_str() -> str:
    return datetime.now().strftime("%Y-%m-%d %H:%M:%S")


def sanitize_excel_string(value: str) -> str:
    if value is None:
        return ""

    value = str(value)
    value = ILLEGAL_CHARACTERS_RE.sub("", value)
    value = value.replace("\r\n", " ")
    value = value.replace("\n", " ")
    value = value.replace("\r", " ")
    value = value.replace("\t", " ")
    value = re.sub(r"\s+", " ", value)
    value = value.strip()
    return value


def normalize_text(value) -> str:
    if value is None:
        return ""
    return sanitize_excel_string(value)


def normalize_for_match(value) -> str:
    if value is None:
        return ""
    return re.sub(r"\s+", "", str(value))


def normalize_int_value(value):
    """
    정수형 데이터 처리
    - 공백, '-', 빈값 => 0
    - 숫자형 문자열 => int 변환
    - 문자열 => 그대로 반환
    """
    if value is None:
        return 0

    text = sanitize_excel_string(value)

    if text in {"", "-"}:
        return 0

    text_no_comma = text.replace(",", "")

    if re.fullmatch(r"[+-]?\d+(\.\d+)?", text_no_comma):
        try:
            return int(float(text_no_comma))
        except ValueError:
            return text

    return text


def normalize_float_value(value):
    """
    실수형 데이터 처리
    - 공백, '-', 빈값 => 0
    - 숫자형 문자열 => float 변환 후 소수점 첫째 자리 반올림
    - 문자열 => 그대로 반환
    """
    if value is None:
        return 0

    text = sanitize_excel_string(value)

    if text in {"", "-"}:
        return 0

    text_no_comma = text.replace(",", "")

    if re.fullmatch(r"[+-]?\d+(\.\d+)?", text_no_comma):
        try:
            return round(float(text_no_comma), 1)
        except ValueError:
            return text

    return text


def extract_part_name(file_path) -> str:
    """
    파트명 추출 우선순위:
    1. FOLDER_PART_MAP: 상위 폴더명 정확 매칭
    2. PART_KEYWORDS: 파일명 내 키워드 검색
    3. 없으면 '-'
    """
    path = Path(file_path) if not isinstance(file_path, Path) else file_path
    parent_folder = path.parent.name
    if parent_folder in FOLDER_PART_MAP:
        return FOLDER_PART_MAP[parent_folder]

    clean_name = sanitize_excel_string(path.name)
    for keyword in PART_KEYWORDS:
        if keyword in clean_name:
            return keyword
    return "-"


def extract_report_stage(file_name: str) -> str:
    clean_name = sanitize_excel_string(file_name)
    for keyword in REPORT_STAGE_KEYWORDS:
        if keyword in clean_name:
            return keyword
    return "-"


def is_temp_file(file_path: Path) -> bool:
    return file_path.name.startswith("~$")


def is_valid_ppt(file_path: Path) -> bool:
    return (
        file_path.is_file()
        and not is_temp_file(file_path)
        and file_path.suffix.lower() in PPT_EXTENSIONS
    )


def compute_file_sha256(file_path: Path) -> str:
    sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        while True:
            chunk = f.read(HASH_CHUNK_SIZE)
            if not chunk:
                break
            sha256.update(chunk)
    return sha256.hexdigest()


def get_file_meta(file_path: Path) -> Dict[str, str]:
    stat = file_path.stat()
    modified_dt = datetime.fromtimestamp(stat.st_mtime)
    created_dt = datetime.fromtimestamp(stat.st_ctime)
    file_size = stat.st_size
    full_path = str(file_path.resolve())
    file_sha256 = compute_file_sha256(file_path)
    part_name = extract_part_name(file_path)  # 폴더명 우선, 파일명 차선
    report_stage = extract_report_stage(file_path.name)

    return {
        "파일식별키": file_sha256,
        "파일명": sanitize_excel_string(file_path.name),
        "전체경로": sanitize_excel_string(full_path),
        "최종수정일시": modified_dt.strftime("%Y-%m-%d %H:%M:%S"),
        "생성일시": created_dt.strftime("%Y-%m-%d %H:%M:%S"),
        "파일크기": str(file_size),
        "파트명": sanitize_excel_string(part_name),
        "보고단계": sanitize_excel_string(report_stage),
    }


# =========================================
# Excel 처리
# =========================================
def load_or_create_workbook(excel_path: Path):
    if excel_path.exists():
        try:
            logger.info(f"기존 엑셀 파일 로드: {excel_path}")
            wb = load_workbook(excel_path)
            return wb
        except Exception as e:
            logger.warning(f"기존 엑셀 파일 로드 실패: {e}")
            logger.warning("손상된 파일로 판단되어 삭제 후 새로 생성합니다.")
            excel_path.unlink()

    logger.info(f"새 엑셀 파일 생성: {excel_path}")
    wb = Workbook()

    if wb.active and wb.active.title == "Sheet":
        wb.active.title = DATA_SHEET_NAME

    return wb


def get_or_create_sheet(wb, sheet_name: str):
    if sheet_name in wb.sheetnames:
        logger.info(f"시트 사용: {sheet_name}")
        return wb[sheet_name]
    logger.info(f"시트 신규 생성: {sheet_name}")
    return wb.create_sheet(sheet_name)


def ensure_history_sheet_if_empty(ws):
    headers = [
        "파일식별키", "파일명", "전체경로", "최종수정일시",
        "파일크기", "처리일시", "처리상태", "메시지"
    ]

    if ws.max_row == 1 and ws.max_column == 1 and ws["A1"].value is None:
        for idx, h in enumerate(headers, start=1):
            ws.cell(row=1, column=idx, value=h)
        logger.info("처리 이력 시트 헤더 신규 생성")
    else:
        logger.info("처리 이력 시트 1행 유지")


def ensure_data_sheet_layout(ws):
    """취합 시트 row 1에 헤더 기록 (항상 덮어씀 — 데이터는 row 2부터)."""
    fixed = ["프로젝트코드", "수행연도", "파트명", "보고단계"]
    # PJ목표/실적/유사/비고 + 사업계획(col4) 순서 — update_summary_sheet 참조 순서 유지
    groups = ["PJ목표", "PJ실적", "PJ유사", "비고", "사업계획"]
    headers = fixed[:]
    for g in groups:
        for m in KPI_METRIC_NAMES:
            headers.append(f"{m}_{g}")
    headers += ["파일명", "최종수정일시", "처리일시"]
    for col_idx, h in enumerate(headers, start=1):
        ws.cell(row=1, column=col_idx, value=h)


def ensure_summary_sheet_layout(ws):
    ws["A1"] = "프로젝트 코드"
    ws["B1"] = "수행연도"
    ws["C1"] = "구분"
    ws["D1"] = "’26년 목표(사업계획)"
    ws["E1"] = "’26년 목표(프로젝트)"
    ws["F1"] = "’26년 실적(프로젝트)"
    ws["G1"] = "’25년 실적(유사)"
    ws["H1"] = "비고"

    ws["A2"] = "E111600126020001"
    ws["B2"] = "2026년"
    ws["C2"] = "교육 만족도 (NPS)"

    ws["A3"] = "E111600126020001"
    ws["B3"] = "2026년"
    ws["C3"] = "그룹 연구개발 2030 전략기술 관련 과정 개발 (과정 건수)"

    ws["A4"] = "E111600126020001"
    ws["B4"] = "2026년"
    ws["C4"] = "그룹 연구개발 2030 전략기술 관련 과정 개발 (교육 내용 구성 적절성)"

    ws["A5"] = "E111600126020001"
    ws["B5"] = "2026년"
    ws["C5"] = "본부별/그룹사별 특화 교육체계 구축 (프로젝트 건수)"

    ws["A6"] = "E111600126020001"
    ws["B6"] = "2026년"
    ws["C6"] = "그룹 내 AI 교육 확대 (고객사 건수)"

    ws["A7"] = "E111600126020001"
    ws["B7"] = "2026년"
    ws["C7"] = "그룹 내 AI 교육 확대 (교육 내용 구성 적절성)"

    ws["A8"] = "E111600126020001"
    ws["B8"] = "2026년"
    ws["C8"] = "정부지원 사업 및 신사업 확대 (매출액, 단위 : 억)"

    ws["A9"] = "E111600126020001"
    ws["B9"] = "2026년"
    ws["C9"] = "정부지원 사업 및 신사업 확대 (신규/기존 사업 건수)"


def append_history(ws, file_meta: Dict[str, str], status: str, message: str):
    next_row = max(ws.max_row + 1, 2)
    ws.cell(next_row, 1, file_meta["파일식별키"])
    ws.cell(next_row, 2, file_meta["파일명"])
    ws.cell(next_row, 3, file_meta["전체경로"])
    ws.cell(next_row, 4, file_meta["최종수정일시"])
    ws.cell(next_row, 5, file_meta["파일크기"])
    ws.cell(next_row, 6, now_str())
    ws.cell(next_row, 7, sanitize_excel_string(status))
    ws.cell(next_row, 8, sanitize_excel_string(message))

    logger.info(f"처리 이력 기록: 상태={status}, 메시지={message}")


# 프로젝트 코드가 이 값이면 "미배정" 상태로 간주 — 서로 다른 프로젝트가 같은 값을
# 공유해도 충돌(덮어쓰기)하지 않도록 중복 판별 키에 파일명을 추가로 사용
PLACEHOLDER_CODES = {"", "0", "생성예정", "미정", "tbd", "(생성 필요)", "선정 시 생성 예정"}
_STAGE_SUFFIXES = ["사전검토", "착수", "중간", "완료", "제안"]
FILENAME_COL = 45  # '취합' 시트 파일명 컬럼(1-based)


def strip_stage_suffix(filename: str) -> str:
    """파일명에서 착수/완료/제안 등 단계 표시를 제거 — 같은 프로젝트의 다른 단계 파일인지
    비교하기 위한 휴리스틱 (완전한 판별은 아니고 충돌 경고용)."""
    name = os.path.splitext(normalize_text(filename))[0]
    for suf in _STAGE_SUFFIXES:
        name = re.sub(rf"[_\[\(]?{re.escape(suf)}[\]\)]?(_수정|_최종)?$", "", name).strip()
    return name


def find_existing_data_row(ws, key1: str, key2: str, key4: str, filename: str = "") -> Optional[int]:
    """프로젝트코드 + 수행연도 + 보고단계 3중 키로 중복 판별 (착수/중간/완료 구분).
    코드가 미배정 플레이스홀더(생성예정 등)면 서로 다른 프로젝트가 같은 값을 공유해도
    충돌하지 않도록 파일명(단계 제거 기준명)도 같이 비교한다."""
    is_placeholder = key1 in PLACEHOLDER_CODES
    file_base = strip_stage_suffix(filename) if is_placeholder else ""
    for row_idx in range(2, ws.max_row + 1):
        v1 = normalize_text(ws.cell(row=row_idx, column=1).value)
        v2 = normalize_text(ws.cell(row=row_idx, column=2).value)
        v4 = normalize_text(ws.cell(row=row_idx, column=4).value)
        if v1 != key1 or v2 != key2 or v4 != key4:
            continue
        if is_placeholder:
            existing_filename = normalize_text(ws.cell(row=row_idx, column=FILENAME_COL).value)
            if strip_stage_suffix(existing_filename) != file_base:
                continue
        return row_idx
    return None


def safe_cell_value(value):
    if isinstance(value, str):
        return sanitize_excel_string(value)
    return value


def cleanup_deleted_files(data_ws, processed_filenames: set) -> int:
    """이번 실행에서 처리된 파일 목록에 없는 행을 취합 시트에서 자동 제거.
    KPI는 전량 재처리 방식이므로 처리된 파일 목록 = 현재 소스 폴더의 전체 파일."""
    rows_to_delete = [
        row_idx
        for row_idx in range(2, data_ws.max_row + 1)
        if normalize_text(data_ws.cell(row=row_idx, column=FILENAME_COL).value) not in processed_filenames
        and normalize_text(data_ws.cell(row=row_idx, column=FILENAME_COL).value) != ""
    ]
    for row_idx in reversed(rows_to_delete):
        data_ws.delete_rows(row_idx)
    if rows_to_delete:
        logger.info(f"[정리] 취합 시트 {len(rows_to_delete)}행 제거 (삭제된 파일)")
    return len(rows_to_delete)


def upsert_data_rows(ws, rows: List[List]) -> Tuple[int, int]:
    inserted = 0
    updated = 0

    for row_data in rows:
        key1 = normalize_text(row_data[0])   # 프로젝트코드
        key2 = normalize_text(row_data[1])   # 수행연도
        key4 = normalize_text(row_data[3])   # 보고단계 (착수/중간/완료 구분)
        new_filename = normalize_text(row_data[44]) if len(row_data) > 44 else ""
        existing_row = find_existing_data_row(ws, key1, key2, key4, new_filename)

        if existing_row is not None:
            existing_filename = normalize_text(ws.cell(row=existing_row, column=FILENAME_COL).value)
            if existing_filename and new_filename and \
               strip_stage_suffix(existing_filename) != strip_stage_suffix(new_filename):
                logger.warning(
                    f"코드 충돌 의심 - 서로 다른 파일이 같은 (코드/연도/단계) 키를 공유함: "
                    f"기존='{existing_filename}' 신규='{new_filename}' 코드={key1} 연도={key2} 단계={key4}"
                )
            target_row = existing_row
            updated += 1
            logger.info(f"기존 데이터 덮어쓰기: 행={target_row}, 코드={key1}, 연도={key2}, 단계={key4}")
        else:
            target_row = max(ws.max_row + 1, 2)
            inserted += 1
            logger.info(f"신규 데이터 추가: 행={target_row}, 코드={key1}, 연도={key2}, 단계={key4}")

        for col_idx, value in enumerate(row_data, start=1):
            ws.cell(row=target_row, column=col_idx, value=safe_cell_value(value))

    return inserted, updated


# =========================================
# KPI 집계 계산
# =========================================
def to_number_or_none(value):
    """
    집계용 숫자 판별:
    - 숫자만 반영
    - N, 공백, 특수문자, 일반 문자열은 제외
    """
    if value is None:
        return None

    if isinstance(value, (int, float)):
        return float(value)

    text = normalize_text(value)

    if text in {"", "-", "N"}:
        return None

    text_no_comma = text.replace(",", "")

    # 순수 숫자형만 허용
    if re.fullmatch(r"[+-]?\d+(\.\d+)?", text_no_comma):
        try:
            return float(text_no_comma)
        except ValueError:
            return None

    return None


def get_numeric_values_from_column(ws, col_idx: int) -> List[float]:
    values = []
    for row_idx in range(2, ws.max_row + 1):
        val = ws.cell(row=row_idx, column=col_idx).value
        num = to_number_or_none(val)
        if num is not None:
            values.append(num)
    return values


def calc_sum(ws, col_idx: int):
    values = get_numeric_values_from_column(ws, col_idx)
    if not values:
        return 0
    return round(sum(values), 2)


def calc_avg(ws, col_idx: int):
    values = get_numeric_values_from_column(ws, col_idx)
    if not values:
        return 0
    return round(sum(values) / len(values), 2)


def calc_new_existing_count_text(ws, col_idx: int) -> str:
    """
    L열 / T열 전용
    - 신규, 기존만 카운트
    - 그 외 값은 무시
    """
    new_count = 0
    existing_count = 0

    for row_idx in range(2, ws.max_row + 1):
        raw = ws.cell(row=row_idx, column=col_idx).value
        text = normalize_text(raw)

        if text == "신규":
            new_count += 1
        elif text == "기존":
            existing_count += 1

    return f"신규:{new_count}건/기존:{existing_count}건"


def update_summary_sheet(summary_ws, data_ws):
    """
    취합 시트 기준으로 kpi 집계 시트 최신화.
    컬럼 매핑:
      E열 = 26년 목표 집계  (취합 col 5~12  : PJ목표)
      F열 = 26년 실적 집계  (취합 col 13~20 : PJ실적)
      G열 = 25년 실적 집계  (취합 col 21~28 : PJ유사)
    신규/기존 형식(col 12, 20, 28)은 별도 카운트 텍스트로 기록.
    """
    ensure_summary_sheet_layout(summary_ws)

    # ── D열: 26년 목표(사업계획, col 37~44) ────────────────
    summary_ws["D2"] = calc_avg(data_ws, 37)
    summary_ws["D3"] = calc_sum(data_ws, 38)
    summary_ws["D4"] = calc_avg(data_ws, 39)
    summary_ws["D5"] = calc_sum(data_ws, 40)
    summary_ws["D6"] = calc_sum(data_ws, 41)
    summary_ws["D7"] = calc_avg(data_ws, 42)
    summary_ws["D8"] = calc_sum(data_ws, 43)
    summary_ws["D9"] = calc_new_existing_count_text(data_ws, 44)

    # ── E열: 26년 목표(프로젝트, PJ목표, col 5~12) ─────────
    summary_ws["E2"] = calc_avg(data_ws, 5)
    summary_ws["E3"] = calc_sum(data_ws, 6)
    summary_ws["E4"] = calc_avg(data_ws, 7)
    summary_ws["E5"] = calc_sum(data_ws, 8)
    summary_ws["E6"] = calc_sum(data_ws, 9)
    summary_ws["E7"] = calc_avg(data_ws, 10)
    summary_ws["E8"] = calc_sum(data_ws, 11)
    summary_ws["E9"] = calc_new_existing_count_text(data_ws, 12)

    # ── F열: 26년 실적 (PJ실적, col 13~20) ────────────────
    summary_ws["F2"] = calc_avg(data_ws, 13)
    summary_ws["F3"] = calc_sum(data_ws, 14)
    summary_ws["F4"] = calc_avg(data_ws, 15)
    summary_ws["F5"] = calc_sum(data_ws, 16)
    summary_ws["F6"] = calc_sum(data_ws, 17)
    summary_ws["F7"] = calc_avg(data_ws, 18)
    summary_ws["F8"] = calc_sum(data_ws, 19)
    summary_ws["F9"] = calc_new_existing_count_text(data_ws, 20)

    # ── G열: 25년 실적 (PJ유사, col 21~28) ────────────────
    summary_ws["G2"] = calc_avg(data_ws, 21)
    summary_ws["G3"] = calc_sum(data_ws, 22)
    summary_ws["G4"] = calc_avg(data_ws, 23)
    summary_ws["G5"] = calc_sum(data_ws, 24)
    summary_ws["G6"] = calc_sum(data_ws, 25)
    summary_ws["G7"] = calc_avg(data_ws, 26)
    summary_ws["G8"] = calc_sum(data_ws, 27)
    summary_ws["G9"] = calc_new_existing_count_text(data_ws, 28)

    logger.info("kpi 집계 시트 최신화 완료 (목표/실적/유사 3개 열)")


# =========================================
# 파일 선택 로직
# =========================================
def list_candidate_ppt_files_with_depth_limit(
    root_dir: Path,
    max_depth: int = MAX_SEARCH_DEPTH,
    exclude_names: Optional[List[str]] = None
) -> List[Path]:
    r"""
    깊이 제한 함수: max_depth까지만 재귀적으로 검색
    """
    exclude_names = exclude_names or []
    files = []

    logger.info(f"PPT 파일 검색 중... (경로: {root_dir}, 최대 깊이: {max_depth})")

    def search_recursive(path: Path, current_depth: int):
        if current_depth > max_depth:
            logger.debug(f"깊이 초과로 검색 중단: {path} (깊이: {current_depth})")
            return

        try:
            for item in path.iterdir():
                if is_valid_ppt(item):
                    if item.name not in exclude_names:
                        files.append(item)
                        logger.debug(f"PPT 파일 발견: {item}")
                elif item.is_dir(follow_symlinks=False):
                    try:
                        search_recursive(item, current_depth + 1)
                    except PermissionError:
                        logger.warning(f"접근 불가: {item}")
                    except Exception as e:
                        logger.warning(f"폴더 검색 오류 {item}: {e}")
        except PermissionError:
            logger.warning(f"접근 불가: {path}")
        except Exception as e:
            logger.warning(f"검색 오류 {path}: {e}")

    search_recursive(root_dir, 0)

    logger.info(f"발견된 PPT 파일: {len(files)}개")
    return files


def get_all_target_files(
    root_dir: Path,
    exclude_names: Optional[List[str]] = None
) -> List[Tuple[Path, Dict[str, str]]]:
    candidates = list_candidate_ppt_files_with_depth_limit(
        root_dir,
        max_depth=MAX_SEARCH_DEPTH,
        exclude_names=exclude_names
    )

    all_files = []

    for file_path in candidates:
        try:
            meta = get_file_meta(file_path)
            all_files.append((file_path, meta))
        except Exception as e:
            logger.exception(f"파일 메타 조회 실패: {file_path} / {e}")

    # 오래된 파일부터 처리 → 최신 파일이 마지막에 upsert되어 완료/중간 실적 데이터가 보존됨
    all_files.sort(
        key=lambda item: (item[0].stat().st_ctime, item[0].stat().st_mtime),
        reverse=False
    )
    return all_files


# =========================================
# PowerPoint 처리
# =========================================
def get_slide_title(slide) -> str:
    try:
        if slide.shapes.title is not None:
            return normalize_text(slide.shapes.title.text)
    except Exception:
        pass
    return ""


def slide_contains_keyword(slide, keyword: str) -> bool:
    normalized_keyword = normalize_for_match(keyword)
    title = get_slide_title(slide)
    if normalized_keyword in normalize_for_match(title):
        return True

    for shape in slide.shapes:
        try:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                shape_text = normalize_text(shape.text)
                if normalized_keyword in normalize_for_match(shape_text):
                    return True
        except Exception:
            continue

    return False


def get_all_tables(slide):
    """슬라이드 내 모든 표를 반환 — 한 슬라이드에 여러 프로젝트 표가 나란히 있는 경우 대응."""
    tables = []
    for idx, shape in enumerate(slide.shapes, start=1):
        try:
            if shape.has_table:
                tables.append(shape.table)
                logger.info(f"표 발견: shape index={idx}")
        except Exception:
            continue

    if not tables:
        logger.warning("슬라이드에서 표를 찾지 못했습니다.")
    return tables


def get_table_text(table, row_1based: int, col_1based: int) -> str:
    try:
        return normalize_text(table.cell(row_1based - 1, col_1based - 1).text)
    except Exception:
        return ""


def extract_record_from_table(
    table,
    source_file_name: str,
    source_modified: str,
    part_name: str,
    report_stage: str
) -> Optional[List]:
    row_count = len(table.rows)
    col_count = len(table.columns)

    logger.info(f"표 크기 확인: rows={row_count}, cols={col_count}")

    if row_count < END_ROW or col_count < max(DATA_COLS):
        logger.warning(
            f"표 구조 부족: 최소 rows={END_ROW}, cols={max(DATA_COLS)} 필요 / 실제 rows={row_count}, cols={col_count}"
        )
        return None

    key1 = get_table_text(table, 2, 1)
    key2 = get_table_text(table, 2, 2)

    logger.info(
        f"추출 키값 확인: 프로젝트 코드='{key1}', 수행연도='{key2}', 파트명='{part_name}', 보고단계='{report_stage}'"
    )

    if not key1 or not key2:
        logger.warning("키값(2행 1열, 2행 2열)이 비어 있어 추출을 건너뜁니다.")
        return None

    row_values = [
        key1,
        key2,
        sanitize_excel_string(part_name),
        sanitize_excel_string(report_stage),
    ]

    for col in DATA_COLS:
        for row in range(START_ROW, END_ROW + 1):
            val = get_table_text(table, row, col)

            if row in FLOAT_ROWS:
                row_values.append(normalize_float_value(val))
            else:
                row_values.append(normalize_int_value(val))

    # col 4: '26년 목표(사업계획) — DATA_COLS 뒤에 추가
    for row in range(START_ROW, END_ROW + 1):
        val = get_table_text(table, row, 4)
        if row in FLOAT_ROWS:
            row_values.append(normalize_float_value(val))
        else:
            row_values.append(normalize_int_value(val))

    row_values.append(sanitize_excel_string(source_file_name))
    row_values.append(source_modified)
    row_values.append(now_str())

    logger.info(f"1개 표 -> 1개 행 변환 완료, 총 컬럼 수={len(row_values)}")
    return row_values


def convert_to_pptx_via_win32(ppt_path: Path) -> Optional[Path]:
    """win32com으로 구버전 .ppt 또는 손상된 .pptx를 새 .pptx로 변환 후 임시 경로 반환."""
    if not WIN32_AVAILABLE:
        logger.error("win32com 없음 — pip install pywin32 필요")
        return None

    tmp_path = None
    ppt_app = None
    prs = None
    try:
        pythoncom.CoInitialize()
        ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
        ppt_app.Visible = 1

        abs_src = str(ppt_path.resolve())
        prs = ppt_app.Presentations.Open(abs_src, True, False, False)

        tmp_fd, tmp_str = tempfile.mkstemp(suffix=".pptx", prefix="kpi_conv_")
        os.close(tmp_fd)
        os.unlink(tmp_str)  # SaveAs가 직접 생성하므로 미리 삭제

        prs.SaveAs(tmp_str, 24)  # 24 = ppSaveAsOpenXMLPresentation
        tmp_path = Path(tmp_str)
        logger.info(f"win32 변환 완료: {ppt_path.name} → {tmp_path}")
        return tmp_path

    except Exception as e:
        logger.error(f"win32 변환 실패: {ppt_path.name} / {e}")
        if tmp_path and tmp_path.exists():
            tmp_path.unlink(missing_ok=True)
        return None
    finally:
        try:
            if prs is not None:
                prs.Close()
        except Exception:
            pass
        try:
            if ppt_app is not None:
                ppt_app.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def _com_cell_text(table_com, row_1based: int, col_1based: int) -> str:
    """COM Table 객체에서 셀 텍스트 추출."""
    try:
        cell = table_com.Cell(row_1based, col_1based)
        return normalize_text(cell.Shape.TextFrame.TextRange.Text)
    except Exception:
        return ""


def _com_slide_has_keyword(slide_com, keyword: str) -> bool:
    """COM Slide 객체에서 키워드 포함 여부 확인."""
    normalized = normalize_for_match(keyword)
    try:
        if slide_com.Shapes.HasTitle:
            title = normalize_text(slide_com.Shapes.Title.TextFrame.TextRange.Text)
            if normalized in normalize_for_match(title):
                return True
    except Exception:
        pass
    try:
        for i in range(1, slide_com.Shapes.Count + 1):
            shp = slide_com.Shapes(i)
            try:
                if shp.HasTextFrame and shp.TextFrame.HasText:
                    if normalized in normalize_for_match(normalize_text(shp.TextFrame.TextRange.Text)):
                        return True
            except Exception:
                continue
    except Exception:
        pass
    return False


def extract_records_from_ppt_via_com(ppt_path: Path, file_meta: Dict[str, str]) -> List[List]:
    """AIP 암호화 등으로 python-pptx + SaveAs가 모두 실패할 때 COM API로 직접 추출."""
    if not WIN32_AVAILABLE:
        raise RuntimeError("win32com 없음")

    ppt_app = None
    prs_com = None
    try:
        pythoncom.CoInitialize()
        ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
        ppt_app.Visible = 1

        prs_com = ppt_app.Presentations.Open(str(ppt_path.resolve()), True, False, False)
        logger.info(f"COM 직접 읽기: {ppt_path.name} (슬라이드 수: {prs_com.Slides.Count})")

        extracted = []
        for s_idx in range(1, prs_com.Slides.Count + 1):
            slide = prs_com.Slides(s_idx)

            if not _com_slide_has_keyword(slide, TITLE_KEYWORD):
                continue

            logger.info(f"COM: 대상 슬라이드 발견 — {s_idx}")

            # 슬라이드 내 모든 표 찾기 — 한 슬라이드에 여러 프로젝트 표가 나란히 있는 경우 대응
            tables_com = []
            for i in range(1, slide.Shapes.Count + 1):
                shp = slide.Shapes(i)
                try:
                    if shp.HasTable:
                        tables_com.append(shp.Table)
                except Exception:
                    continue

            if not tables_com:
                logger.warning(f"COM: 슬라이드 {s_idx} 표 없음")
                continue

            for table_com in tables_com:
                rows_n = table_com.Rows.Count
                cols_n = table_com.Columns.Count
                if rows_n < END_ROW or cols_n < max(DATA_COLS):
                    logger.warning(f"COM: 표 구조 부족 rows={rows_n} cols={cols_n}")
                    continue

                key1 = _com_cell_text(table_com, 2, KEY_COL_1)
                key2 = _com_cell_text(table_com, 2, KEY_COL_2)
                if not key1 or not key2:
                    continue

                row_values = [key1, key2,
                              sanitize_excel_string(file_meta["파트명"]),
                              sanitize_excel_string(file_meta["보고단계"])]

                for col in DATA_COLS:
                    for row in range(START_ROW, END_ROW + 1):
                        val = _com_cell_text(table_com, row, col)
                        if row in FLOAT_ROWS:
                            row_values.append(normalize_float_value(val))
                        else:
                            row_values.append(normalize_int_value(val))

                # col 4: '26년 목표(사업계획)
                for row in range(START_ROW, END_ROW + 1):
                    val = _com_cell_text(table_com, row, 4)
                    if row in FLOAT_ROWS:
                        row_values.append(normalize_float_value(val))
                    else:
                        row_values.append(normalize_int_value(val))

                row_values.append(sanitize_excel_string(file_meta["파일명"]))
                row_values.append(file_meta["최종수정일시"])
                row_values.append(now_str())

                extracted.append(row_values)
                logger.info(f"COM: 데이터 1건 추출 완료")

        logger.info(f"COM 추출 총 {len(extracted)}건: {ppt_path.name}")
        return extracted

    finally:
        try:
            if prs_com is not None:
                prs_com.Close()
        except Exception:
            pass
        try:
            if ppt_app is not None:
                ppt_app.Quit()
        except Exception:
            pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


def extract_records_from_ppt(ppt_path: Path, file_meta: Dict[str, str]) -> List[List]:
    converted_tmp: Optional[Path] = None
    actual_path = ppt_path

    # .ppt 파일은 바로 변환
    if ppt_path.suffix.lower() == ".ppt":
        logger.info(f".ppt 감지 → win32 변환 시도: {ppt_path.name}")
        converted_tmp = convert_to_pptx_via_win32(ppt_path)
        if converted_tmp is None:
            raise ValueError(f".ppt 변환 실패 (win32com 필요): {ppt_path.name}")
        actual_path = converted_tmp

    try:
        try:
            prs = Presentation(str(actual_path))
        except Exception as e:
            # 1차: win32 변환 후 재시도
            if converted_tmp is None:
                logger.warning(f"pptx 파싱 실패 → win32 변환 재시도: {ppt_path.name} ({e})")
                converted_tmp = convert_to_pptx_via_win32(ppt_path)
                if converted_tmp is not None:
                    actual_path = converted_tmp
                    try:
                        prs = Presentation(str(actual_path))
                    except Exception as e2:
                        # 2차: 변환된 파일도 실패(AIP 재암호화) → COM 직접 추출
                        logger.warning(f"변환 후도 실패 → COM 직접 추출: {ppt_path.name} ({e2})")
                        return extract_records_from_ppt_via_com(ppt_path, file_meta)
                else:
                    # SaveAs 자체가 실패한 경우(AIP 쓰기 제한) → COM 직접 추출
                    logger.warning(f"win32 변환 실패 → COM 직접 추출: {ppt_path.name}")
                    return extract_records_from_ppt_via_com(ppt_path, file_meta)
            else:
                # 이미 변환된 상태인데도 실패 → COM 직접 추출
                logger.warning(f"변환 파일 파싱 실패 → COM 직접 추출: {ppt_path.name} ({e})")
                return extract_records_from_ppt_via_com(ppt_path, file_meta)

        extracted = []
        logger.info(f"PowerPoint 열기: {ppt_path.name} (슬라이드 수: {len(prs.slides)})")

        for slide_idx, slide in enumerate(prs.slides, start=1):
            title_text = get_slide_title(slide)
            logger.info(f"슬라이드 {slide_idx} 제목: {title_text}")

            if not slide_contains_keyword(slide, TITLE_KEYWORD):
                logger.info(f"슬라이드 {slide_idx}: KPI 키워드 불일치")
                continue

            logger.info(f"대상 슬라이드 발견: {slide_idx}")

            tables = get_all_tables(slide)
            if not tables:
                logger.warning(f"슬라이드 {slide_idx}: 표를 찾지 못했습니다.")
                continue

            for table_idx, table in enumerate(tables, start=1):
                record = extract_record_from_table(
                    table=table,
                    source_file_name=file_meta["파일명"],
                    source_modified=file_meta["최종수정일시"],
                    part_name=file_meta["파트명"],
                    report_stage=file_meta["보고단계"]
                )

                if record:
                    extracted.append(record)
                    logger.info(f"슬라이드 {slide_idx} 표 {table_idx}: 데이터 1건 추출")
                else:
                    logger.warning(f"슬라이드 {slide_idx} 표 {table_idx}: 표 구조 또는 키값 부족으로 건너뜀")

        logger.info(f"총 추출 건수: {len(extracted)}")
        return extracted

    finally:
        if converted_tmp and converted_tmp.exists():
            converted_tmp.unlink(missing_ok=True)


# =========================================
# 파일 1건 처리
# =========================================
def process_single_file(
    ppt_path: Path,
    file_meta: Dict[str, str],
    data_ws,
    history_ws
) -> Tuple[bool, str, List[List]]:
    try:
        records = extract_records_from_ppt(ppt_path, file_meta)

        if not records:
            msg = "조건에 맞는 슬라이드 또는 표 데이터를 찾지 못했습니다."
            append_history(history_ws, file_meta, "처리완료", msg)
            logger.info(msg)
            return True, msg, []

        msg = f"추출 {len(records)}건"
        append_history(history_ws, file_meta, "처리완료", msg)

        logger.info(f"파일 처리 완료: {ppt_path.name} / {msg}")
        return True, msg, records

    except Exception as e:
        msg = str(e)
        append_history(history_ws, file_meta, "실패", msg)
        logger.exception(f"파일 처리 중 예외 발생: {ppt_path.name}")
        return False, msg, []


# =========================================
# 메인 처리
# =========================================
def main():
    logger.info("=== KPI PowerPoint 데이터 추출 작업 시작 ===")
    logger.info(f"대상 디렉토리: {ROOT_DIR}")
    logger.info(f"검색 깊이: {MAX_SEARCH_DEPTH}단계")
    logger.info("처리 방식: 순차 처리(안정성 우선)")
    logger.info("엑셀 1행 컬럼값은 현재 상태 유지")
    logger.info("파일명 추출 항목: 파트명(C열), 보고단계(D열)")
    logger.info("kpi 집계 시트는 추출 시마다 최신화 업데이트")
    logger.info("집계 규칙: L/T열 제외 숫자만 반영, L/T열은 신규/기존만 카운트")

    if not ROOT_DIR.exists():
        logger.error(f"대상 폴더가 존재하지 않습니다: {ROOT_DIR}")
        print(f"[ERROR] 대상 폴더가 존재하지 않습니다: {ROOT_DIR}")
        return

    excel_path = _DATA_DIR / TARGET_EXCEL_NAME

    wb = load_or_create_workbook(excel_path)
    data_ws = get_or_create_sheet(wb, DATA_SHEET_NAME)
    history_ws = get_or_create_sheet(wb, HISTORY_SHEET_NAME)
    summary_ws = get_or_create_sheet(wb, SUMMARY_SHEET_NAME)

    ensure_data_sheet_layout(data_ws)
    ensure_history_sheet_if_empty(history_ws)
    ensure_summary_sheet_layout(summary_ws)

    # --retry: kpi_aip_failed.txt에 기록된 파일만 재처리
    if RETRY_MODE:
        if not AIP_FAILED_FILE.exists():
            logger.info("[retry] 실패 목록 파일이 없습니다. 먼저 일반 실행으로 추출하세요.")
            print("[retry] kpi_aip_failed.txt 없음. 일반 실행 먼저 하세요.")
            return
        failed_paths = [p.strip() for p in AIP_FAILED_FILE.read_text(encoding="utf-8").splitlines() if p.strip()]
        all_target_files = [
            (Path(p), {
                '파트명': '',
                '보고단계': '',
                '전체경로': p,
                '파일명': Path(p).name,
                '수정일시': '',
                '파일크기': '',
                '해시': '',
            })
            for p in failed_paths if Path(p).exists()
        ]
        logger.info(f"[retry] AIP 실패 목록 {len(all_target_files)}개 재처리 시작")
        _retry_remaining = [str(p) for p, _ in all_target_files]
        # 목록은 성공 시에만 제거 — 재실패 파일은 유지
    else:
        all_target_files = get_all_target_files(
            ROOT_DIR,
            exclude_names=[TARGET_EXCEL_NAME]
        )

    logger.info(f"전체 대상 파일 수: {len(all_target_files)}")

    if not all_target_files:
        update_summary_sheet(summary_ws, data_ws)
        wb.save(excel_path)
        logger.info("처리할 PowerPoint 파일이 없습니다.")
        print("[INFO] 처리할 PowerPoint 파일이 없습니다.")
        return

    total_success = 0
    total_fail = 0
    all_records = []

    for ppt_path, file_meta in all_target_files:
        logger.info(f"전체 추출 요청: {ppt_path.name}")
        logger.info(f"파일 경로: {file_meta['전체경로']}")
        logger.info(f"추출된 파트명: {file_meta['파트명']}")
        logger.info(f"추출된 보고단계: {file_meta['보고단계']}")

        success, message, records = process_single_file(
            ppt_path=ppt_path,
            file_meta=file_meta,
            data_ws=data_ws,
            history_ws=history_ws
        )

        if records:
            all_records.extend(records)

        if success:
            total_success += 1
            if RETRY_MODE:
                _retry_remaining = [p for p in _retry_remaining if p != str(ppt_path)]
        else:
            total_fail += 1
            # AIP 관련 실패 시 retry 목록에 기록 (일반 모드)
            if not RETRY_MODE and any(kw in message for kw in ["AIP", "암호화", "PackageNotFound", "win32"]):
                with open(AIP_FAILED_FILE, "a", encoding="utf-8") as _f:
                    _f.write(str(ppt_path) + "\n")
                logger.info(f"[retry 대상 기록] {ppt_path.name}")

        wb.save(excel_path)

    if all_records:
        inserted, updated = upsert_data_rows(data_ws, all_records)
        logger.info(f"전체 데이터 적재 완료: 신규 {inserted}건 / 덮어쓰기 {updated}건")

    # 소스 폴더에 없는 파일 데이터 자동 정리
    if not RETRY_MODE:
        processed_filenames = {meta["파일명"] for _, meta in all_target_files}
        cleanup_deleted_files(data_ws, processed_filenames)

    update_summary_sheet(summary_ws, data_ws)

    wb.save(excel_path)

    # retry 모드: 최종 실패 파일 목록 다시 저장
    if RETRY_MODE:
        AIP_FAILED_FILE.write_text("\n".join(_retry_remaining) + ("\n" if _retry_remaining else ""), encoding="utf-8")
        if _retry_remaining:
            logger.info(f"[retry] 재실패 {len(_retry_remaining)}개 → {AIP_FAILED_FILE} 에 유지")
        else:
            logger.info("[retry] 모든 파일 처리 완료 — 실패 목록 초기화")

    logger.info(f"전체 처리 결과: 성공 {total_success}건 / 실패 {total_fail}건")
    logger.info(f"총 데이터 건수: {len(all_records)}건")
    logger.info(f"저장 파일: {excel_path}")
    logger.info("=== KPI PowerPoint 데이터 추출 작업 종료 ===")

    print("[DONE] 작업 완료")
    print(f" - 처리 성공 파일 수: {total_success}")
    print(f" - 처리 실패 파일 수: {total_fail}")
    print(f" - 총 데이터 건수: {len(all_records)}")
    print(f" - 저장 파일: {excel_path}")
    print(f" - 로그 파일: {_DATA_DIR / LOG_FILE_NAME}")


if __name__ == "__main__":
    main()
