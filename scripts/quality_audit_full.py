import sys, io, os, re
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import extract_financial_ppt as ex
from pptx import Presentation

BASE_DIR = ex.BASE_DIR
TITLE_KEYWORD = ex.TITLE_KEYWORD
FIN_HINTS = ["프로젝트 코드", "총 매출", "지출", "직접", "공통원가", "경상", "이익"]

results = []

for root, _, files in os.walk(BASE_DIR):
    for name in files:
        if name.startswith("~$") or not name.lower().endswith(".pptx"):
            continue
        full = os.path.join(root, name)
        if ex.is_aip_encrypted(full):
            continue  # AIP는 별도 처리(이미 파악됨)
        try:
            prs = Presentation(full)
        except Exception as e:
            results.append({"file": full, "issue": "열기실패", "detail": str(e)})
            continue

        matched_by_title = False
        found_alt_table = False
        note_header_variant = None

        for slide in prs.slides:
            has_kw = any(sh.has_text_frame and TITLE_KEYWORD in sh.text_frame.text for sh in slide.shapes)
            for sh in slide.shapes:
                if not sh.has_table:
                    continue
                tbl = sh.table
                try:
                    headers = [tbl.cell(0, c).text.strip() for c in range(len(tbl.columns))]
                except Exception:
                    continue
                hits = sum(1 for h in FIN_HINTS if any(h in cell for cell in headers))
                if hits < 3:
                    continue
                if has_kw:
                    matched_by_title = True
                    last_header = headers[-1] if headers else ""
                    if last_header and "비고" not in last_header:
                        note_header_variant = last_header
                else:
                    found_alt_table = True

        if not matched_by_title and found_alt_table:
            results.append({"file": full, "issue": "구템플릿(제목 불일치)", "detail": ""})
        elif not matched_by_title and not found_alt_table:
            results.append({"file": full, "issue": "재무데이터 없음", "detail": ""})
        elif matched_by_title and note_header_variant:
            results.append({"file": full, "issue": "비고헤더텍스트상이", "detail": note_header_variant})

for r in results:
    print(f"{r['issue']}\t{os.path.basename(r['file'])}\t{r['detail']}")
