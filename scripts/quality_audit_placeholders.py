import sys, io, os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

import extract_financial_ppt as ex
from pptx import Presentation
from collections import defaultdict

PLACEHOLDERS = {'생성예정', '0', '미정', '(생성 필요)', '선정 시 생성 예정', 'tbd', ''}

BASE_DIR = ex.BASE_DIR
TITLE_KEYWORD = ex.TITLE_KEYWORD

rows_by_placeholder = defaultdict(list)

for root, _, files in os.walk(BASE_DIR):
    for name in files:
        if name.startswith("~$"):
            continue
        if not name.lower().endswith(".pptx"):
            continue
        full = os.path.join(root, name)
        if ex.is_aip_encrypted(full):
            continue
        try:
            prs = Presentation(full)
        except Exception:
            continue
        for slide in prs.slides:
            if not any(sh.has_text_frame and TITLE_KEYWORD in sh.text_frame.text for sh in slide.shapes):
                continue
            best = None
            bc = 0
            for sh in slide.shapes:
                if sh.has_table and len(sh.table.columns) > bc:
                    bc = len(sh.table.columns)
                    best = sh.table
            if best is None or bc < 10:
                continue
            for r in range(1, len(best.rows)):
                try:
                    code = ex.safe_str(best.cell(r, 0).text)
                    gubun = ex.safe_str(best.cell(r, 1).text)
                except Exception:
                    continue
                if code in PLACEHOLDERS:
                    rows_by_placeholder[code].append((os.path.basename(full), gubun))

for code, items in rows_by_placeholder.items():
    print(f'코드="{code}" 사용 건수: {len(items)}')
    for fn, gubun in items:
        print(f'   [{gubun}] {fn}')
